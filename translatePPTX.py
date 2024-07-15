import os
import argparse
import requests
from pptx import Presentation
from tqdm import tqdm

'''
You need to pip install to your python enviroment the following command:

pip install requests python-pptx tqdm
'''

# Set your Google Cloud API key
API_KEY = ""

def translate_text(text, target_language):
    url = f"https://translation.googleapis.com/language/translate/v2?key={API_KEY}"
    headers = {
        "Content-Type": "application/json"
    }
    body = {
        "q": text,
        "target": target_language,
        "format": "text"
    }
    response = requests.post(url, headers=headers, json=body)
    if response.status_code == 200:
        return response.json()['data']['translations'][0]['translatedText']
    else:
        print(f"Error translating text: {response.status_code} {response.text}")
        return text

def translate_shape_text(shape, target_language):
    if not hasattr(shape, "text_frame") or not shape.text_frame:
        return

    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            translated_text = translate_text(run.text, target_language)
            run.text = translated_text

def process_presentation(input_file, target_language):
    print(f"Opening {input_file}")
    try:
        input_ppt = Presentation(input_file)
    except Exception as e:
        print(f"Error opening file {input_file}: {e}")
        return

    slide_count = len(input_ppt.slides)
    
    with tqdm(total=slide_count, desc="Translating", unit="slide") as pbar:
        for i, slide in enumerate(input_ppt.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    try:
                        translate_shape_text(shape, target_language)
                    except Exception as e:
                        print(f"Error processing shape on slide {i}: {e}")
            pbar.update(1)

    output_file = f"{target_language}_{input_file}"
    try:
        input_ppt.save(output_file)
        print(f"\nSaved as {output_file}")
    except Exception as e:
        print(f"Error saving file {output_file}: {e}")

def main():
    # for link output to console (this sometimes works)
    def link(uri, label=None):
        if label is None: 
            label = uri
        parameters = ''

        # OSC 8 ; params ; URI ST <name> OSC 8 ;; ST 
        escape_mask = '\033]8;{};{}\033\\{}\033]8;;\033\\'

        return escape_mask.format(parameters, uri, label)
    # ------------------------

    print("Please use the 2 letter language code for the argument!")
    print("Example language syntax:")
    print("English: en")
    print("Spanish: es")
    print("French: fr")
    print("German: de")
    print("")
    print("See Full List of ISO 639 Languages here: " + link('https://en.wikipedia.org/wiki/List_of_ISO_639_language_codes'))
    print("")

    parser = argparse.ArgumentParser(description="Translate a PowerPoint presentation. Usage: python3 translatePPTX.py <input_pptx_file> <target_language>")
    parser.add_argument("input_file", help="Path to the input PowerPoint file")
    parser.add_argument("target_language", help="Target language for translation (ex: 'en' for English, 'es' for Spanish)")
    args = parser.parse_args()

    # check if users code is valid before running it (it will cause a bunch of errors @ api endpoint)
    valid_language_codes = [
        'aa', 'ab', 'ae', 'af', 'ak', 'am', 'an', 'ar', 'as', 'av', 'ay', 'az', 'ba', 'be', 'bg', 'bi', 'bm', 'bn', 'bo', 
        'br', 'bs', 'ca', 'ce', 'ch', 'co', 'cr', 'cs', 'cu', 'cv', 'cy', 'da', 'de', 'dv', 'dz', 'ee', 'el', 'en', 'eo', 
        'es', 'et', 'eu', 'fa', 'ff', 'fi', 'fj', 'fo', 'fr', 'fy', 'ga', 'gd', 'gl', 'gn', 'gu', 'gv', 'ha', 'he', 'hi', 
        'ho', 'hr', 'ht', 'hu', 'hy', 'hz', 'ia', 'id', 'ie', 'ig', 'ii', 'ik', 'io', 'is', 'it', 'iu', 'ja', 'jv', 'ka', 
        'kg', 'ki', 'kj', 'kk', 'kl', 'km', 'kn', 'ko', 'kr', 'ks', 'ku', 'kv', 'kw', 'ky', 'la', 'lb', 'lg', 'li', 'ln', 
        'lo', 'lt', 'lu', 'lv', 'mg', 'mh', 'mi', 'mk', 'ml', 'mn', 'mr', 'ms', 'mt', 'my', 'na', 'nb', 'nd', 'ne', 'ng', 
        'nl', 'nn', 'no', 'nr', 'nv', 'ny', 'oc', 'oj', 'om', 'or', 'os', 'pa', 'pi', 'pl', 'ps', 'pt', 'qu', 'rm', 'rn', 
        'ro', 'ru', 'rw', 'sa', 'sc', 'sd', 'se', 'sg', 'si', 'sk', 'sl', 'sm', 'sn', 'so', 'sq', 'sr', 'ss', 'st', 'su', 
        'sv', 'sw', 'ta', 'te', 'tg', 'th', 'ti', 'tk', 'tl', 'tn', 'to', 'tr', 'ts', 'tt', 'tw', 'ty', 'ug', 'uk', 'ur', 
        'uz', 've', 'vi', 'vo', 'wa', 'wo', 'xh', 'yi', 'yo', 'za', 'zh', 'zu'
    ]

    if args.target_language in valid_language_codes:
        process_presentation(args.input_file, args.target_language)
    else:
        print("ERROR: NOT A VALID LANGUAGE CODE")
        print("")
        print("ERROR: Please submit a valid language code")
        print("")

        return -1


if __name__ == "__main__":
    main()
